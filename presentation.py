import io
from pptx import Presentation
from pptx.util import Inches

class PresentationGenerator:
    def create_presentation(self, slide_data):
        presentation = Presentation()
        for slide_info in slide_data:
            slide_layout = presentation.slide_layouts[slide_info['layout']]
            slide = presentation.slides.add_slide(slide_layout)
            title = slide.shapes.title
            title.text = slide_info['heading']
            body = slide.shapes.placeholders[1].text_frame
            body.text = slide_info['body']

        presentation_file = io.BytesIO()
        presentation.save(presentation_file)
        presentation_file.seek(0)
        return presentation_file
