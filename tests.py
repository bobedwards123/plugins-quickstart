print('running tests for the project')

# Path: main.py
from utils import create_presentation_new
from pptx import Presentation
import io
from pptx.util import Pt

dictionary = {
  "slide_data": [
    {
      "heading": "Introduction to Dogs",
      "body": "Dogs, also known as Canis lupus familiaris, are domesticated mammals that are commonly kept as pets. They are known for their loyalty and companionship.",
      "layout": 2
    },
    {
      "heading": "History of Dogs",
      "body": "Dogs were likely the first domesticated animals, with evidence suggesting domestication as early as 15,000 years ago. They have been bred over generations for various behaviors, sensory capabilities, and physical attributes.",
      "layout": 2
    },
     {
      "heading": "History of Dogs",
      "body": "Dogs were likely the first domesticated animals, with evidence suggesting domestication as early as 15,000 years ago. They have been bred over generations for various behaviors, sensory capabilities, and physical attributes.",
      "layout": 2
    },
     {
      "heading": "History of Dogs",
      "body": "Dogs were likely the first domesticated animals, with evidence suggesting domestication as early as 15,000 years ago. They have been bred over generations for various behaviors, sensory capabilities, and physical attributes.",
      "layout": 3
    },
    {
      "heading": "History of Dogs",
      "body": "Dogs were likely the first domesticated animals, with evidence suggesting domestication as early as 15,000 years ago. They have been bred over generations for various behaviors, sensory capabilities, and physical attributes.",
      "layout": 3
    },
    {
      "heading": "History of Dogs",
      "body": "Dogs were likely the first domesticated animals, with evidence suggesting domestication as early as 15,000 years ago. They have been bred over generations for various behaviors, sensory capabilities, and physical attributes.",
      "layout": 3
    },
    {
      "heading": "History of Dogs",
      "body": "Dogs were likely the first domesticated animals, with evidence suggesting domestication as early as 15,000 years ago. They have been bred over generations for various behaviors, sensory capabilities, and physical attributes.",
      "layout": 3
    },
    {
      "heading": "History of Dogs",
      "body": "Dogs were likely the first domesticated animals, with evidence suggesting domestication as early as 15,000 years ago. They have been bred over generations for various behaviors, sensory capabilities, and physical attributes.",
      "layout": 3
    },
    {
      "heading": "History of Dogs",
      "body": "Dogs were likely the first domesticated animals, with evidence suggesting domestication as early as 15,000 years ago. They have been bred over generations for various behaviors, sensory capabilities, and physical attributes.",
      "layout": 3
    },
    {
      "heading": "History of Dogs",
      "body": "Dogs were likely the first domesticated animals, with evidence suggesting domestication as early as 15,000 years ago. They have been bred over generations for various behaviors, sensory capabilities, and physical attributes.",
      "layout": 3
    },
    {
      "heading": "History of Dogs",
      "body": "Dogs were likely the first domesticated animals, with evidence suggesting domestication as early as 15,000 years ago. They have been bred over generations for various behaviors, sensory capabilities, and physical attributes.",
      "layout": 3
    },
    {
      "heading": "History of Dogs",
      "body": "Dogs were likely the first domesticated animals, with evidence suggesting domestication as early as 15,000 years ago. They have been bred over generations for various behaviors, sensory capabilities, and physical attributes.",
      "layout": 3
    },
    {
      "heading": "History of Dogs",
      "body": "Dogs were likely the first domesticated animals, with evidence suggesting domestication as early as 15,000 years ago. They have been bred over generations for various behaviors, sensory capabilities, and physical attributes.",
      "layout": 3
    },
    {
      "heading": "History of Dogs",
      "body": "Dogs were likely the first domesticated animals, with evidence suggesting domestication as early as 15,000 years ago. They have been bred over generations for various behaviors, sensory capabilities, and physical attributes.",
      "layout": 3
    },
    {
      "heading": "History of Dogs",
      "body": "Dogs were likely the first domesticated animals, with evidence suggesting domestication as early as 15,000 years ago. They have been bred over generations for various behaviors, sensory capabilities, and physical attributes.",
      "layout": 3
    },
    {
      "heading": "History of Dogs",
      "body": "Dogs were likely the first domesticated animals, with evidence suggesting domestication as early as 15,000 years ago. They have been bred over generations for various behaviors, sensory capabilities, and physical attributes.",
      "layout": 3
    },
    {
      "heading": "History of Dogs",
      "body": "Dogs were likely the first domesticated animals, with evidence suggesting domestication as early as 15,000 years ago. They have been bred over generations for various behaviors, sensory capabilities, and physical attributes.",
      "layout": 3
    }
  ]
}


# it should return a file path
# print("it should create a presentation")
# output_link = create_presentation(dictionary["slide_data"])
# print(output_link)

# # assert that the file path is correct with .pptx at the end
# assert output_link[-5:] == ".pptx"

# # it should delete a presentation
# print("it should delete a presentation")
# delete_presentation(output_link)


# print("it should be able to take in a template as an argument and then use that template to create the presentation")
# # it should be able to take in a template as an argument and then use that template to create the presentation
# output_link = create_presentation(dictionary["slide_data"], template="x_template.pptx")

# it should create a presentation from a template when provided

# output_link = create_presentation(dictionary["slide_data"], template="x_template.pptx")
# print(output_link)
# print("it should create a presentation from a template when provided")
new_schema = {
  "slide_data": [
      {
      "heading": "Introduction to Dogs",
      "content": [
          "Dogs, also known as Canis lupus familiaris, are domesticated mammals that are commonly kept as pets",
          "They are known for their loyalty and companionship."
      ]
      },
           {
      "heading": "Introduction to Dogs",
      "content": [
          "Dogs, also known as Canis lupus familiaris, are domesticated mammals that are commonly kept as pets",
          "They are known for their loyalty and companionship."
      ]
      }
  ]
}

# def create_presentation_new(slide_data):
#     presentation = Presentation()
#     index = 0
#     for slide in slide_data:
        
#         heading = slide.get("heading", "")
#         content = slide.get("content", [])
#         layout = None
#         if index == 0:
#           layout = presentation.slide_layouts[2]  
#         else:
#           layout = presentation.slide_layouts[1]
#         slide = presentation.slides.add_slide(layout)
#         title = slide.shapes.placeholders[0].text_frame
#         title.text = heading
        
#         content_holder = slide.shapes.placeholders[1].text_frame

#         for paragraph in content:
#             content_holder.add_paragraph().text = paragraph

#     # create a filepath for the presentation
#     filepath = "./presentations/presentation.pptx"
#     presentation.save(filepath)


create_presentation_new(new_schema)

    