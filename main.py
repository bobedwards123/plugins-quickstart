import json
import quart
import quart_cors
from quart import request, send_file, jsonify, send_from_directory
from presentation import PresentationGenerator
from utils import create_presentation_new
import uuid
import os
import io
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches

app = quart_cors.cors(quart.Quart(__name__), allow_origin="https://chat.openai.com")

# Directory to store generated presentations
PRESENTATION_DIR = "./presentations"
TEMPLATE_DIR = "./templates"



def create_presentation(slide_data, template=None):
    # Generate a unique identifier for this presentation
    presentation_id = str(uuid.uuid4())
    # Create a file path for the presentation
    presentation_path = os.path.join(PRESENTATION_DIR, f"{presentation_id}.pptx")

    if template:
        template_link = os.path.join(TEMPLATE_DIR, template)
        presentation = Presentation(template_link)
    else:
        presentation = Presentation()

    for slide_dict in slide_data:
        # Iterate through slide layouts until we find one with title and content placeholders
        for layout in presentation.slide_layouts:
            placeholders = layout.placeholders
            if any(ph.placeholder_format.idx == 0 for ph in placeholders) and \
               any(ph.placeholder_format.idx == 1 for ph in placeholders):
                slide_layout = layout
                break
        else:
            print("Error: No suitable slide layout found")
            continue

        slide = presentation.slides.add_slide(slide_layout)
        title = slide.placeholders[0]
        content = slide.placeholders[1]
        title.text = slide_dict['heading']
        content.text = slide_dict['body']

    presentation.save(presentation_path)
    return presentation_path

def delete_presentation(presentation_id):
    # Create a file path for the presentation
    presentation_path = presentation_id
    # Delete the file
    os.remove(presentation_path)



@app.post("/presentation/link_new")
async def presentation_link_new():
    slide_data = (await request.get_json())['slide_data']
    presentation_file = create_presentation_new(slide_data, PRESENTATION_DIR)
    print(presentation_file)
    
    # Generate a unique download link for the file
    download_link = f"http://{request.host}/presentation/download/{os.path.basename(presentation_file)}"
    return jsonify({"download_link": download_link})



@app.post('/presentation/link')
async def presentation_link():
    slide_data = (await request.get_json())['slide_data']
    presentation_file = create_presentation(slide_data)
    print(presentation_file)
    
    # Generate a unique download link for the file
    download_link = f"http://{request.host}/presentation/download/{os.path.basename(presentation_file)}"
    return jsonify({"download_link": download_link})

@app.get('/presentation/download/<presentation_id>')
async def presentation_download(presentation_id):
    # Serve the file from the presentations directory
    return await send_from_directory(PRESENTATION_DIR, presentation_id, as_attachment=True)

# create and endpoint to get the template
@app.get('/templates/<template_id>')
async def template_download(template_id):
    # Serve the file from the presentations directory
    return await send_from_directory(TEMPLATE_DIR, template_id, as_attachment=True)

@app.get("/logo.png")
async def plugin_logo():
    filename = 'logo.png'
    return await quart.send_file(filename, mimetype='image/png')

@app.get("/.well-known/ai-plugin.json")
async def plugin_manifest():
    host = request.headers['Host']
    with open("./.well-known/ai-plugin.json") as f:
        text = f.read()
        return quart.Response(text, mimetype="text/json")

@app.get("/openapi.yaml")
async def openapi_spec():
    host = request.headers['Host']
    with open("openapi.yaml") as f:
        text = f.read()
        return quart.Response(text, mimetype="text/yaml")

@app.post('/presentation')
async def presentation_endpoint():
    slide_data = (await request.get_json())['slide_data']
    presentation_generator = PresentationGenerator()
    presentation_file = presentation_generator.create_presentation(slide_data)
    return await send_file(presentation_file, as_attachment=True, attachment_filename='presentation.pptx')

# @app.post('/presentation/download')
# async def presentation_download():
#     slide_data = (await request.get_json())['slide_data']
#     presentation_generator = PresentationGenerator()
#     presentation_file = presentation_generator.create_presentation(slide_data)
#     return await send_file(presentation_file, as_attachment=True, attachment_filename='presentation.pptx')

# @app.post('/presentation/link')
# async def presentation_link():
#     slide_data = (await request.get_json())['slide_data']
#     presentation_generator = PresentationGenerator()
#     presentation_file = presentation_generator.create_presentation(slide_data)
#     # Generate a unique download link for the file
#     download_link = f"http://{request.host}/presentation/download"
#     return jsonify({"download_link": download_link})

def main():
    app.run(debug=True, host="0.0.0.0", port=5003)

if __name__ == "__main__":
    main()
