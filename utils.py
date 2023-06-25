from pptx import Presentation
import uuid
import os

def create_presentation_new(slide_data, PRESENTATION_DIR):
        # Generate a unique identifier for this presentation
    presentation_id = str(uuid.uuid4())
    # Create a file path for the presentation
    presentation_path = os.path.join(PRESENTATION_DIR, f"{presentation_id}.pptx")
    presentation = Presentation()
    # slide_data = input_data["slide_data"]
    index = 0
    for slide in slide_data:
        
        heading = slide.get("heading", "")
        content = slide.get("content", [])
        layout = None
        if index == 0:
          layout = presentation.slide_layouts[2]  
        else:
          layout = presentation.slide_layouts[1]
        slide = presentation.slides.add_slide(layout)
        title = slide.shapes.placeholders[0].text_frame
        title.text = heading
        
        content_holder = slide.shapes.placeholders[1].text_frame

        for paragraph in content:
            content_holder.add_paragraph().text = paragraph
        index += 1

    presentation.save(presentation_path)
    return presentation_path